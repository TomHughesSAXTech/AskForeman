"""
Enhanced Document Converter for Azure Function
Supports: PDF, Word, Excel, PowerPoint, Text, CSV, and more
"""

import io
import base64
import logging
import chardet
from typing import Optional, Dict, Any, List
import json

# PDF handling
try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

# Word document handling
try:
    from docx import Document
except ImportError:
    Document = None

# Excel handling
try:
    import openpyxl
    from openpyxl import load_workbook
except ImportError:
    openpyxl = None

# PowerPoint handling
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# CSV handling
import csv

logger = logging.getLogger(__name__)


class DocumentConverter:
    """Universal document converter for multiple file types"""
    
    def __init__(self):
        self.supported_types = self._get_supported_types()
    
    def _get_supported_types(self) -> Dict[str, bool]:
        """Check which file types are supported based on available libraries"""
        return {
            'pdf': PdfReader is not None,
            'docx': Document is not None,
            'xlsx': openpyxl is not None,
            'pptx': Presentation is not None,
            'txt': True,
            'csv': True,
            'md': True,
            'json': True,
        }
    
    def extract_text(self, content: bytes, file_name: str, mime_type: str = None) -> str:
        """
        Extract text from document based on file type
        
        Args:
            content: File content as bytes
            file_name: Name of the file
            mime_type: MIME type of the file (optional)
        
        Returns:
            Extracted text as string
        """
        # Determine file type from extension
        file_ext = file_name.lower().split('.')[-1] if '.' in file_name else ''
        
        # Route to appropriate extractor
        extractors = {
            'pdf': self._extract_pdf_text,
            'docx': self._extract_docx_text,
            'doc': self._extract_docx_text,  # Try docx extractor for .doc
            'xlsx': self._extract_xlsx_text,
            'xls': self._extract_xlsx_text,  # Try xlsx extractor for .xls
            'pptx': self._extract_pptx_text,
            'ppt': self._extract_pptx_text,  # Try pptx extractor for .ppt
            'txt': self._extract_plain_text,
            'md': self._extract_plain_text,
            'csv': self._extract_csv_text,
            'json': self._extract_json_text,
            'xml': self._extract_plain_text,
            'html': self._extract_plain_text,
            'htm': self._extract_plain_text,
        }
        
        extractor = extractors.get(file_ext, self._extract_plain_text)
        
        try:
            text = extractor(content, file_name)
            logger.info(f"Successfully extracted text from {file_name} ({file_ext})")
            return text
        except Exception as e:
            logger.error(f"Failed to extract text from {file_name}: {str(e)}")
            # Fallback to plain text extraction
            return self._extract_plain_text(content, file_name)
    
    def _extract_pdf_text(self, content: bytes, file_name: str) -> str:
        """Extract text from PDF files"""
        if not PdfReader:
            logger.warning("PyPDF2 not installed, falling back to plain text")
            return self._extract_plain_text(content, file_name)
        
        try:
            pdf_file = io.BytesIO(content)
            pdf_reader = PdfReader(pdf_file)
            
            text_parts = []
            for page_num, page in enumerate(pdf_reader.pages, 1):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"--- Page {page_num} ---\n{page_text}")
                except Exception as e:
                    logger.warning(f"Failed to extract page {page_num}: {str(e)}")
            
            return "\n\n".join(text_parts) if text_parts else "No text content found in PDF"
        except Exception as e:
            logger.error(f"PDF extraction failed: {str(e)}")
            raise
    
    def _extract_docx_text(self, content: bytes, file_name: str) -> str:
        """Extract text from Word documents"""
        if not Document:
            logger.warning("python-docx not installed, falling back to plain text")
            return self._extract_plain_text(content, file_name)
        
        try:
            doc_file = io.BytesIO(content)
            doc = Document(doc_file)
            
            text_parts = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Extract text from tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        table_text.append(" | ".join(row_text))
                if table_text:
                    text_parts.append("\n".join(table_text))
            
            return "\n\n".join(text_parts) if text_parts else "No text content found in document"
        except Exception as e:
            logger.error(f"Word document extraction failed: {str(e)}")
            raise
    
    def _extract_xlsx_text(self, content: bytes, file_name: str) -> str:
        """Extract text from Excel files"""
        if not openpyxl:
            logger.warning("openpyxl not installed, falling back to plain text")
            return self._extract_plain_text(content, file_name)
        
        try:
            excel_file = io.BytesIO(content)
            workbook = load_workbook(excel_file, read_only=True, data_only=True)
            
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_data = []
                
                # Extract data from cells
                for row in sheet.iter_rows(values_only=True):
                    row_data = []
                    for cell in row:
                        if cell is not None:
                            # Convert to string and clean
                            cell_str = str(cell).strip()
                            if cell_str:
                                row_data.append(cell_str)
                    
                    if row_data:
                        sheet_data.append(" | ".join(row_data))
                
                if sheet_data:
                    text_parts.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(sheet_data))
            
            workbook.close()
            return "\n\n".join(text_parts) if text_parts else "No text content found in spreadsheet"
        except Exception as e:
            logger.error(f"Excel extraction failed: {str(e)}")
            raise
    
    def _extract_pptx_text(self, content: bytes, file_name: str) -> str:
        """Extract text from PowerPoint presentations"""
        if not Presentation:
            logger.warning("python-pptx not installed, falling back to plain text")
            return self._extract_plain_text(content, file_name)
        
        try:
            pptx_file = io.BytesIO(content)
            presentation = Presentation(pptx_file)
            
            text_parts = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text.strip())
                    
                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))
                
                if slide_text:
                    text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
            
            return "\n\n".join(text_parts) if text_parts else "No text content found in presentation"
        except Exception as e:
            logger.error(f"PowerPoint extraction failed: {str(e)}")
            raise
    
    def _extract_csv_text(self, content: bytes, file_name: str) -> str:
        """Extract text from CSV files"""
        try:
            # Detect encoding
            encoding = self._detect_encoding(content)
            text_content = content.decode(encoding)
            
            # Parse CSV
            csv_reader = csv.reader(io.StringIO(text_content))
            
            rows = []
            for row in csv_reader:
                if row:
                    rows.append(" | ".join(row))
            
            return "\n".join(rows) if rows else "No content found in CSV"
        except Exception as e:
            logger.error(f"CSV extraction failed: {str(e)}")
            return self._extract_plain_text(content, file_name)
    
    def _extract_json_text(self, content: bytes, file_name: str) -> str:
        """Extract text from JSON files"""
        try:
            # Detect encoding and decode
            encoding = self._detect_encoding(content)
            text_content = content.decode(encoding)
            
            # Parse JSON and format it
            json_data = json.loads(text_content)
            
            # Convert to formatted string
            formatted = json.dumps(json_data, indent=2, ensure_ascii=False)
            
            # Extract meaningful text values
            text_parts = []
            
            def extract_values(obj, prefix=""):
                if isinstance(obj, dict):
                    for key, value in obj.items():
                        if isinstance(value, (str, int, float, bool)):
                            text_parts.append(f"{prefix}{key}: {value}")
                        elif isinstance(value, (dict, list)):
                            extract_values(value, f"{prefix}{key}.")
                elif isinstance(obj, list):
                    for i, item in enumerate(obj):
                        if isinstance(item, (str, int, float, bool)):
                            text_parts.append(f"{prefix}[{i}]: {item}")
                        else:
                            extract_values(item, f"{prefix}[{i}].")
            
            extract_values(json_data)
            
            return "\n".join(text_parts) if text_parts else formatted
        except Exception as e:
            logger.error(f"JSON extraction failed: {str(e)}")
            return self._extract_plain_text(content, file_name)
    
    def _extract_plain_text(self, content: bytes, file_name: str) -> str:
        """Extract plain text with encoding detection"""
        try:
            # Try to detect encoding
            encoding = self._detect_encoding(content)
            text = content.decode(encoding)
            return text if text.strip() else "No text content found"
        except Exception as e:
            logger.error(f"Plain text extraction failed: {str(e)}")
            # Last resort: try common encodings
            for enc in ['utf-8', 'latin-1', 'cp1252', 'ascii']:
                try:
                    return content.decode(enc)
                except:
                    continue
            return "Unable to decode file content"
    
    def _detect_encoding(self, content: bytes) -> str:
        """Detect the encoding of byte content"""
        try:
            result = chardet.detect(content)
            return result['encoding'] or 'utf-8'
        except:
            return 'utf-8'


# Main function to be used by Azure Function
def extract_document_text(content: bytes, file_name: str, mime_type: str = None) -> str:
    """
    Main entry point for document text extraction
    
    Args:
        content: File content as bytes
        file_name: Name of the file
        mime_type: MIME type of the file (optional)
    
    Returns:
        Extracted text as string
    """
    converter = DocumentConverter()
    return converter.extract_text(content, file_name, mime_type)


# Compatibility function for existing code
def extract_pdf_text(content: bytes) -> str:
    """Legacy function for PDF extraction"""
    converter = DocumentConverter()
    return converter._extract_pdf_text(content, "document.pdf")
